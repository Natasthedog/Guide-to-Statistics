from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from difflib import SequenceMatcher

import pandas as pd
from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from deck.engine.time_period import CompanyWeekMapper, _coerce_yearwk, _find_company_week_value
from deck.engine.pptx.text import _replace_placeholders_in_slide_runs

from .waterfall_service_layer import WaterfallSlideMapper

logger = logging.getLogger(__name__)

_MARKER_PHRASE = "Media KPIs Summary"
_NET_REVENUE_ROI_TITLE = "Net Revenue ROI"
_MEDIA_INVESTMENT_TITLE = "Media Investment"
_MEDIA_INCREMENTAL_VOLUME_TITLE = "Media Incremental Volume (Kg)"
_BRAND_PLACEHOLDER = "<BRAND>"
_MAT_1_PLACEHOLDER = "<MAT-1>"
_MAT_PLACEHOLDER = "<MAT>"
_NR_ROI_YOY_PLACEHOLDER = "<NR ROI YoY%>"

_BRAND_COLUMN_CANDIDATES = ("Brand", "Target Brand", "Brand Name")
_EFFECT_TYPE_COLUMN_CANDIDATES = ("Effect Type", "EffectType")
_MAT_COLUMN_CANDIDATES = ("MAT", "Mat")
_PROFIT_COLUMN_CANDIDATES = (
    "Profit_Incremental_Total Adstock",
    "Profit Incremental Total Adstock",
    "Profit Incremental TotalAdstock",
)
_SPEND_COLUMN_CANDIDATES = ("Spend", "Total Spend")
_VOLUME_INCREMENTAL_COLUMN_CANDIDATES = (
    "Pumped up Volume_Incremental_Total Adstock YoY",
    "Pumped up Volume Incremental Total Adstock YoY",
    "Pumped Up Volume_Incremental_Total Adstock YoY",
    "PumpedUpVolumeIncrementalTotalAdstockYoY",
)
_ALLOWED_EFFECT_TYPE_TOKENS = {"target", "targethalo"}

_CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_NS = {"c": _CHART_NS}

_GREEN = RGBColor(0, 176, 80)
_RED = RGBColor(192, 0, 0)
_GRAY = RGBColor(127, 127, 127)
_WHITE = RGBColor(255, 255, 255)

_NR_ROI_TEXT_OVERLAY_NAME = "NR ROI YoY Overlay"


def _normalize_token(value: object) -> str:
    if value is None or pd.isna(value):
        return ""
    return "".join(ch for ch in str(value).strip().lower() if ch.isalnum())


def _is_blank(value: object) -> bool:
    if value is None:
        return True
    if isinstance(value, str) and not value.strip():
        return True
    return bool(pd.isna(value)) if not isinstance(value, str) else False


def _iter_shapes_recursive(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shape.shapes)


def _set_shape_text_preserve_formatting(shape, text: str, *, rgb: RGBColor | None = None) -> None:
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        paragraph = text_frame.add_paragraph()
    else:
        paragraph = text_frame.paragraphs[0]

    if not paragraph.runs:
        run = paragraph.add_run()
    else:
        run = paragraph.runs[0]

    run.text = str(text)
    if rgb is not None:
        run.font.color.rgb = rgb

    for extra_run in paragraph.runs[1:]:
        extra_run.text = ""
    for extra_paragraph in text_frame.paragraphs[1:]:
        for extra_run in extra_paragraph.runs:
            extra_run.text = ""


def _remove_shape(shape) -> None:
    try:
        element = shape._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)
    except Exception:
        pass


def _remove_named_shapes(slide, shape_name: str) -> None:
    target = _normalize_token(shape_name)
    for shape in list(_iter_shapes_recursive(slide.shapes)):
        if _normalize_token(getattr(shape, "name", "")) == target:
            _remove_shape(shape)


def _apply_overlay_font_from_source(run, source_shape, *, default_size_pt: float = 9.3) -> None:
    source_run = None
    try:
        if getattr(source_shape, "has_text_frame", False):
            for paragraph in source_shape.text_frame.paragraphs:
                for candidate_run in paragraph.runs:
                    source_run = candidate_run
                    break
                if source_run is not None:
                    break
    except Exception:
        source_run = None

    try:
        if source_run is not None and source_run.font.name:
            run.font.name = source_run.font.name
    except Exception:
        pass
    try:
        if source_run is not None and source_run.font.size is not None:
            run.font.size = source_run.font.size
        else:
            run.font.size = Pt(default_size_pt)
    except Exception:
        pass
    try:
        if source_run is not None and source_run.font.bold is not None:
            run.font.bold = source_run.font.bold
    except Exception:
        pass
    try:
        if source_run is not None and source_run.font.italic is not None:
            run.font.italic = source_run.font.italic
    except Exception:
        pass


def _add_centered_text_overlay(slide, source_shape, text: str, *, rgb: RGBColor = _WHITE, name: str = _NR_ROI_TEXT_OVERLAY_NAME):
    textbox = slide.shapes.add_textbox(source_shape.left, source_shape.top, source_shape.width, source_shape.height)
    try:
        textbox.name = name
    except Exception:
        pass
    try:
        textbox.fill.background()
    except Exception:
        pass
    try:
        textbox.line.fill.background()
    except Exception:
        pass
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    try:
        from pptx.enum.text import PP_ALIGN
        paragraph.alignment = PP_ALIGN.CENTER
    except Exception:
        pass
    run = paragraph.add_run()
    run.text = str(text)
    run.font.color.rgb = rgb
    _apply_overlay_font_from_source(run, source_shape)
    return textbox


def _find_media_kpis_template_slide(prs: Presentation, marker_phrase: str = _MARKER_PHRASE):
    marker = str(marker_phrase or "").strip()
    if not marker:
        return None
    marker_norm = _normalize_token(marker)

    for slide in prs.slides:
        for shape in _iter_shapes_recursive(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text_value = shape.text_frame.text or ""
            if marker in text_value:
                return slide
            if marker_norm and marker_norm in _normalize_token(text_value):
                return slide
    return None


def _column_match_score(value: object, candidates: tuple[str, ...]) -> float:
    token = _normalize_token(value)
    if not token:
        return 0.0

    candidate_tokens = [_normalize_token(candidate) for candidate in candidates if _normalize_token(candidate)]
    if token in candidate_tokens:
        return 1.0

    best_score = 0.0
    for candidate in candidate_tokens:
        score = SequenceMatcher(None, token, candidate).ratio()
        if candidate in token or token in candidate:
            score = max(score, 0.95)
        best_score = max(best_score, score)
    return best_score


def _normalize_media_template_brand_df(media_template_brand_df: pd.DataFrame) -> pd.DataFrame:
    if media_template_brand_df is None or media_template_brand_df.empty:
        raise ValueError("Media_Template_Brand upload is empty.")

    df = media_template_brand_df.copy()
    if isinstance(df.columns, pd.MultiIndex):
        flattened_columns = []
        for column in df.columns:
            parts = [
                str(part).strip()
                for part in column
                if str(part).strip() and not str(part).strip().lower().startswith("unnamed")
            ]
            flattened_columns.append(parts[-1] if parts else "")
        df.columns = flattened_columns

    best_column_score = max((_column_match_score(column, _BRAND_COLUMN_CANDIDATES) for column in df.columns), default=0.0)
    if best_column_score >= 0.82:
        return df

    if df.empty:
        return df

    first_row = df.iloc[0].tolist()
    best_row_score = max((_column_match_score(value, _BRAND_COLUMN_CANDIDATES) for value in first_row), default=0.0)
    if best_row_score < 0.82:
        return df

    normalized_headers = []
    fallback_index = 1
    for value in first_row:
        header = str(value).strip() if value is not None and not pd.isna(value) else ""
        if not header:
            header = f"Column{fallback_index}"
        normalized_headers.append(header)
        fallback_index += 1

    normalized_df = df.iloc[1:].copy().reset_index(drop=True)
    normalized_df.columns = normalized_headers
    return normalized_df


def _resolve_column(df: pd.DataFrame, candidates: tuple[str, ...], *, threshold: float = 0.82) -> str:
    exact_matches: list[str] = []
    scored: list[tuple[float, str]] = []

    for column in df.columns:
        original = str(column)
        score = _column_match_score(column, candidates)
        if score >= 1.0:
            exact_matches.append(original)
            continue
        if score > 0:
            scored.append((score, original))

    if exact_matches:
        unique = list(dict.fromkeys(exact_matches))
        if len(unique) == 1:
            return unique[0]
        raise ValueError(f"Ambiguous columns for {candidates!r}: {unique!r}")

    if not scored:
        raise ValueError(f"Could not resolve column from candidates {candidates!r}.")

    scored.sort(key=lambda item: item[0], reverse=True)
    best_score, best_column = scored[0]
    if best_score < threshold:
        raise ValueError(
            "Could not reliably resolve column from candidates "
            f"{candidates!r}; best score {best_score:.2f} for column {best_column!r}."
        )
    return best_column


@dataclass
class ScopeYearRange:
    start_year: int
    end_year: int


@dataclass
class RoiValues:
    mat_1_roi: float | None
    mat_roi: float | None


@dataclass
class MediaInvestmentValues:
    mat_1_spend: float | None
    mat_spend: float | None


@dataclass
class MediaIncrementalVolumeValues:
    mat_1_volume: float | None
    mat_volume: float | None


def _ordered_unique_brands(media_template_brand_df: pd.DataFrame) -> list[str]:
    normalized_df = _normalize_media_template_brand_df(media_template_brand_df)
    brand_column = _resolve_column(normalized_df, _BRAND_COLUMN_CANDIDATES)
    unique_brands: list[str] = []
    seen: set[str] = set()

    for value in normalized_df[brand_column].tolist():
        if value is None or pd.isna(value):
            continue
        brand = str(value).strip()
        if not brand:
            continue
        key = _normalize_token(brand)
        if not key or key in seen:
            continue
        seen.add(key)
        unique_brands.append(brand)
    return unique_brands


def _resolve_scope_year_range(scope_df: pd.DataFrame) -> ScopeYearRange:
    if scope_df is None or scope_df.empty:
        raise ValueError("Scope dataframe is required to resolve <MAT-1>/<MAT> years.")

    try:
        start_week = _find_company_week_value(scope_df, "First week of modelling")
        end_week = _find_company_week_value(scope_df, "Last week of modelling")
        start_year = CompanyWeekMapper._yearwk_to_monday(_coerce_yearwk(start_week)).year
        end_year = CompanyWeekMapper._yearwk_to_monday(_coerce_yearwk(end_week)).year
        return ScopeYearRange(start_year=start_year, end_year=end_year)
    except Exception:
        pass

    rows = list(scope_df.itertuples(index=False))
    start_year: int | None = None
    end_year: int | None = None

    for row in rows:
        label = _normalize_token(row[0] if len(row) > 0 else "")
        value = row[1] if len(row) > 1 else None
        if value is None or pd.isna(value):
            continue
        try:
            year_value = int(float(value))
        except Exception:
            continue

        if label in {"startyear", "mat1", "year1", "startmatyear", "matminus1"} and start_year is None:
            start_year = year_value
        if label in {"endyear", "mat", "year2", "endmatyear"} and end_year is None:
            end_year = year_value
        if label == "year" and end_year is None:
            end_year = year_value

    if end_year is not None and start_year is None:
        start_year = end_year - 1

    if start_year is None or end_year is None:
        raise ValueError("Could not resolve start/end year from scope file for Media KPIs Summary slide.")

    return ScopeYearRange(start_year=start_year, end_year=end_year)


def _apply_media_kpis_placeholders(slide, *, brand: str, start_year: int, end_year: int) -> bool:
    replacements = {
        _BRAND_PLACEHOLDER: str(brand),
        _MAT_1_PLACEHOLDER: str(start_year),
        _MAT_PLACEHOLDER: str(end_year),
    }
    return _replace_placeholders_in_slide_runs(slide, replacements)


def _coerce_numeric(value: object) -> float:
    if value is None or pd.isna(value):
        return 0.0
    try:
        return float(value)
    except Exception:
        try:
            return float(str(value).strip())
        except Exception:
            return 0.0


def _mat_bucket_key(value: object, *, year_range: ScopeYearRange) -> str:
    token = _normalize_token(value)
    if token in {"year1", "mat1", "matminus1", str(year_range.start_year)}:
        return "year1"
    if token in {"year2", "mat2", "mat", str(year_range.end_year)}:
        return "year2"
    return ""


def _compute_brand_roi_values(
    media_template_brand_df: pd.DataFrame,
    *,
    brand: str,
    year_range: ScopeYearRange,
) -> RoiValues:
    df = _normalize_media_template_brand_df(media_template_brand_df)

    brand_column = _resolve_column(df, _BRAND_COLUMN_CANDIDATES)
    effect_type_column = _resolve_column(df, _EFFECT_TYPE_COLUMN_CANDIDATES)
    mat_column = _resolve_column(df, _MAT_COLUMN_CANDIDATES)
    profit_column = _resolve_column(df, _PROFIT_COLUMN_CANDIDATES)
    spend_column = _resolve_column(df, _SPEND_COLUMN_CANDIDATES)

    brand_key = _normalize_token(brand)
    brand_series = df[brand_column].map(_normalize_token)
    effect_series = df[effect_type_column].map(_normalize_token)
    mat_series = df[mat_column].map(lambda value: _mat_bucket_key(value, year_range=year_range))

    def _roi_for_bucket(bucket_key: str) -> float | None:
        mask = (
            brand_series.eq(brand_key)
            & effect_series.isin(_ALLOWED_EFFECT_TYPE_TOKENS)
            & mat_series.eq(bucket_key)
        )
        subset = df.loc[mask]
        if subset.empty:
            return 0.0

        numerator = pd.to_numeric(subset[profit_column], errors="coerce").fillna(0.0).sum()
        denominator = pd.to_numeric(subset[spend_column], errors="coerce").fillna(0.0).sum()
        if denominator == 0:
            return 0.0
        return round(float(numerator) / float(denominator), 2)

    return RoiValues(
        mat_1_roi=_roi_for_bucket("year1"),
        mat_roi=_roi_for_bucket("year2"),
    )


def _compute_brand_media_investment_values(
    media_template_brand_df: pd.DataFrame,
    *,
    brand: str,
    year_range: ScopeYearRange,
) -> MediaInvestmentValues:
    df = _normalize_media_template_brand_df(media_template_brand_df)

    brand_column = _resolve_column(df, _BRAND_COLUMN_CANDIDATES)
    effect_type_column = _resolve_column(df, _EFFECT_TYPE_COLUMN_CANDIDATES)
    mat_column = _resolve_column(df, _MAT_COLUMN_CANDIDATES)
    spend_column = _resolve_column(df, _SPEND_COLUMN_CANDIDATES)

    brand_key = _normalize_token(brand)
    brand_series = df[brand_column].map(_normalize_token)
    effect_series = df[effect_type_column].map(_normalize_token)
    mat_series = df[mat_column].map(lambda value: _mat_bucket_key(value, year_range=year_range))

    def _spend_for_bucket(bucket_key: str) -> float | None:
        mask = (
            brand_series.eq(brand_key)
            & effect_series.isin(_ALLOWED_EFFECT_TYPE_TOKENS)
            & mat_series.eq(bucket_key)
        )
        subset = df.loc[mask]
        if subset.empty:
            return 0.0
        return float(pd.to_numeric(subset[spend_column], errors="coerce").fillna(0.0).sum())

    return MediaInvestmentValues(
        mat_1_spend=_spend_for_bucket("year1"),
        mat_spend=_spend_for_bucket("year2"),
    )


def _compute_brand_media_incremental_volume_values(
    media_template_brand_df: pd.DataFrame,
    *,
    brand: str,
    year_range: ScopeYearRange,
) -> MediaIncrementalVolumeValues:
    df = _normalize_media_template_brand_df(media_template_brand_df)

    brand_column = _resolve_column(df, _BRAND_COLUMN_CANDIDATES)
    effect_type_column = _resolve_column(df, _EFFECT_TYPE_COLUMN_CANDIDATES)
    mat_column = _resolve_column(df, _MAT_COLUMN_CANDIDATES)
    volume_column = _resolve_column(df, _VOLUME_INCREMENTAL_COLUMN_CANDIDATES)

    brand_key = _normalize_token(brand)
    brand_series = df[brand_column].map(_normalize_token)
    effect_series = df[effect_type_column].map(_normalize_token)
    mat_series = df[mat_column].map(lambda value: _mat_bucket_key(value, year_range=year_range))

    def _volume_for_bucket(bucket_key: str) -> float | None:
        mask = (
            brand_series.eq(brand_key)
            & effect_series.isin(_ALLOWED_EFFECT_TYPE_TOKENS)
            & mat_series.eq(bucket_key)
        )
        subset = df.loc[mask]
        if subset.empty:
            return 0.0
        return float(pd.to_numeric(subset[volume_column], errors="coerce").fillna(0.0).sum())

    return MediaIncrementalVolumeValues(
        mat_1_volume=_volume_for_bucket("year1"),
        mat_volume=_volume_for_bucket("year2"),
    )


def _find_chart_shape_by_title(slide, chart_title: str = _NET_REVENUE_ROI_TITLE):
    target = str(chart_title or "").strip()
    if not target:
        return None
    target_norm = _normalize_token(target)

    for shape in _iter_shapes_recursive(slide.shapes):
        if not getattr(shape, "has_chart", False):
            continue
        try:
            chart = shape.chart
            title_text = chart.chart_title.text_frame.text if chart.has_title else ""
        except Exception:
            title_text = ""
        if target in title_text:
            return shape
        if target_norm and target_norm in _normalize_token(title_text):
            return shape
    return None


def _load_chart_workbook(chart):
    xlsx_blob = chart.part.chart_workbook.xlsx_part.blob
    return load_workbook(io.BytesIO(xlsx_blob))


def _save_chart_workbook(chart, workbook) -> None:
    stream = io.BytesIO()
    workbook.save(stream)
    chart.part.chart_workbook.xlsx_part.blob = stream.getvalue()


def _parse_formula(formula: str) -> tuple[str, str]:
    if "!" not in formula:
        raise ValueError(f"Unsupported chart formula: {formula!r}")
    sheet_name, cell_range = formula.split("!", 1)
    sheet_name = sheet_name.strip()
    if sheet_name.startswith("'") and sheet_name.endswith("'"):
        sheet_name = sheet_name[1:-1].replace("''", "'")
    return sheet_name, cell_range.replace("$", "")


def _cells_from_formula(workbook, formula: str):
    sheet_name, cell_range = _parse_formula(formula)
    ws = workbook[sheet_name]
    if ":" in cell_range:
        rows = ws[cell_range]
        return ws, [cell for row in rows for cell in row]
    return ws, [ws[cell_range]]


def _values_from_formula(workbook, formula: str) -> list[object]:
    _ws, cells = _cells_from_formula(workbook, formula)
    return [cell.value for cell in cells]


def _clear_children(node) -> None:
    for child in list(node):
        node.remove(child)


def _ensure_child(parent, tag_local: str):
    child = parent.find(f"c:{tag_local}", namespaces=_NS)
    if child is None:
        child = etree.SubElement(parent, f"{{{_CHART_NS}}}{tag_local}")
    return child


def _update_str_cache(cache_node, values: list[object]) -> None:
    _clear_children(cache_node)
    pt_count = etree.SubElement(cache_node, f"{{{_CHART_NS}}}ptCount")
    pt_count.set("val", str(len(values)))
    for idx, value in enumerate(values):
        pt = etree.SubElement(cache_node, f"{{{_CHART_NS}}}pt")
        pt.set("idx", str(idx))
        v = etree.SubElement(pt, f"{{{_CHART_NS}}}v")
        v.text = "" if value is None else str(value)


def _update_num_cache(cache_node, values: list[object]) -> None:
    _clear_children(cache_node)
    pt_count = etree.SubElement(cache_node, f"{{{_CHART_NS}}}ptCount")
    pt_count.set("val", str(len(values)))
    for idx, value in enumerate(values):
        pt = etree.SubElement(cache_node, f"{{{_CHART_NS}}}pt")
        pt.set("idx", str(idx))
        v = etree.SubElement(pt, f"{{{_CHART_NS}}}v")
        if value is None or _is_blank(value):
            v.text = ""
        else:
            numeric = _coerce_numeric(value)
            v.text = format(numeric, "g")


def _update_series_and_category_caches(chart, workbook) -> None:
    root = chart.part._element
    for ser in root.findall(".//c:ser", namespaces=_NS):
        tx = ser.find("c:tx", namespaces=_NS)
        if tx is not None:
            str_ref = tx.find("c:strRef", namespaces=_NS)
            if str_ref is not None:
                f_node = str_ref.find("c:f", namespaces=_NS)
                if f_node is not None and f_node.text:
                    values = _values_from_formula(workbook, f_node.text)
                    cache = _ensure_child(str_ref, "strCache")
                    _update_str_cache(cache, values)
            else:
                literal_v = tx.find("c:v", namespaces=_NS)
                if literal_v is not None and literal_v.text is not None:
                    literal_v.text = str(literal_v.text)

        cat = ser.find("c:cat", namespaces=_NS)
        if cat is not None:
            str_ref = cat.find("c:strRef", namespaces=_NS)
            num_ref = cat.find("c:numRef", namespaces=_NS)
            if str_ref is not None:
                f_node = str_ref.find("c:f", namespaces=_NS)
                if f_node is not None and f_node.text:
                    values = _values_from_formula(workbook, f_node.text)
                    cache = _ensure_child(str_ref, "strCache")
                    _update_str_cache(cache, values)
            elif num_ref is not None:
                f_node = num_ref.find("c:f", namespaces=_NS)
                if f_node is not None and f_node.text:
                    values = _values_from_formula(workbook, f_node.text)
                    cache = _ensure_child(num_ref, "numCache")
                    _update_num_cache(cache, values)

        num_ref = ser.find("c:val/c:numRef", namespaces=_NS)
        if num_ref is not None:
            f_node = num_ref.find("c:f", namespaces=_NS)
            if f_node is not None and f_node.text:
                values = _values_from_formula(workbook, f_node.text)
                cache = _ensure_child(num_ref, "numCache")
                _update_num_cache(cache, values)


def _find_workbook_header_column(ws, targets: tuple[str, ...]) -> int | None:
    target_tokens = {_normalize_token(value) for value in targets if _normalize_token(value)}
    for col_idx in range(1, ws.max_column + 1):
        token = _normalize_token(ws.cell(row=1, column=col_idx).value)
        if token and token in target_tokens:
            return col_idx
    return None


def _ensure_total_row(ws) -> int:
    for row_idx in range(2, ws.max_row + 1):
        if _normalize_token(ws.cell(row=row_idx, column=1).value) == "total":
            return row_idx
    if ws.max_row < 2:
        ws.cell(row=2, column=1, value="Total")
        return 2
    return 2


def _update_two_value_chart_workbook(
    slide,
    *,
    chart_title: str,
    year_range: ScopeYearRange,
    mat_1_value: float | None,
    mat_value: float | None,
    error_label: str,
) -> None:
    chart_shape = _find_chart_shape_by_title(slide, chart_title)
    if chart_shape is None:
        logger.info("No '%s' chart found on Media KPIs Summary slide.", chart_title)
        return

    chart = chart_shape.chart
    workbook = _load_chart_workbook(chart)
    ws = workbook.active

    start_col = _find_workbook_header_column(ws, (_MAT_1_PLACEHOLDER, str(year_range.start_year)))
    end_col = _find_workbook_header_column(ws, (_MAT_PLACEHOLDER, str(year_range.end_year)))
    total_row = _ensure_total_row(ws)

    if start_col is None or end_col is None:
        raise ValueError(
            f"Could not find {error_label} chart year columns in embedded workbook. "
            f"Expected headers like {_MAT_1_PLACEHOLDER!r}/{_MAT_PLACEHOLDER!r} or {year_range.start_year}/{year_range.end_year}."
        )

    # Leave Column1 untouched. Only inject year headers and the two metric values.
    ws.cell(row=1, column=start_col, value=str(year_range.start_year))
    ws.cell(row=1, column=end_col, value=str(year_range.end_year))
    ws.cell(row=total_row, column=start_col, value=mat_1_value)
    ws.cell(row=total_row, column=end_col, value=mat_value)

    _save_chart_workbook(chart, workbook)
    _update_series_and_category_caches(chart, workbook)



def _update_net_revenue_roi_chart(
    slide,
    *,
    brand: str,
    year_range: ScopeYearRange,
    media_template_brand_df: pd.DataFrame,
) -> RoiValues:
    roi_values = _compute_brand_roi_values(
        media_template_brand_df,
        brand=brand,
        year_range=year_range,
    )

    _update_two_value_chart_workbook(
        slide,
        chart_title=_NET_REVENUE_ROI_TITLE,
        year_range=year_range,
        mat_1_value=roi_values.mat_1_roi,
        mat_value=roi_values.mat_roi,
        error_label="ROI",
    )
    return roi_values



def _update_media_investment_chart(
    slide,
    *,
    brand: str,
    year_range: ScopeYearRange,
    media_template_brand_df: pd.DataFrame,
) -> MediaInvestmentValues:
    media_investment_values = _compute_brand_media_investment_values(
        media_template_brand_df,
        brand=brand,
        year_range=year_range,
    )

    _update_two_value_chart_workbook(
        slide,
        chart_title=_MEDIA_INVESTMENT_TITLE,
        year_range=year_range,
        mat_1_value=media_investment_values.mat_1_spend,
        mat_value=media_investment_values.mat_spend,
        error_label="Media Investment",
    )
    return media_investment_values


def _update_media_incremental_volume_chart(
    slide,
    *,
    brand: str,
    year_range: ScopeYearRange,
    media_template_brand_df: pd.DataFrame,
) -> MediaIncrementalVolumeValues:
    media_incremental_volume_values = _compute_brand_media_incremental_volume_values(
        media_template_brand_df,
        brand=brand,
        year_range=year_range,
    )

    _update_two_value_chart_workbook(
        slide,
        chart_title=_MEDIA_INCREMENTAL_VOLUME_TITLE,
        year_range=year_range,
        mat_1_value=media_incremental_volume_values.mat_1_volume,
        mat_value=media_incremental_volume_values.mat_volume,
        error_label="Media Incremental Volume",
    )
    return media_incremental_volume_values


def _find_nr_roi_arrow_shape(slide):
    for shape in _iter_shapes_recursive(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text_value = shape.text_frame.text or ""
        if _NR_ROI_YOY_PLACEHOLDER in text_value:
            return shape
    return None


def _nr_roi_change_parts(previous_value: float, current_value: float) -> tuple[str, RGBColor, float, str]:
    previous_value = float(previous_value or 0.0)
    current_value = float(current_value or 0.0)

    if abs(current_value - previous_value) < 1e-12:
        return "=", _GRAY, 0.0, "flat"

    if abs(previous_value) < 1e-12:
        if current_value > 0:
            return "—", _GREEN, 0.0, "up"
        return "=", _GRAY, 0.0, "flat"

    change_value = (current_value - previous_value) / previous_value
    pct_text = f"{abs(change_value) * 100:.0f}%"

    if change_value > 0:
        return pct_text, _GREEN, 0.0, "up"
    if change_value < 0:
        return pct_text, _RED, 180.0, "down"
    return "=", _GRAY, 0.0, "flat"


def _update_nr_roi_yoy_arrow(slide, *, mat_1_roi: float, mat_roi: float) -> None:
    arrow_shape = _find_nr_roi_arrow_shape(slide)
    if arrow_shape is None:
        logger.info("No NR ROI YoY arrow shape found on Media KPIs Summary slide.")
        return

    _remove_named_shapes(slide, _NR_ROI_TEXT_OVERLAY_NAME)
    change_text, change_rgb, arrow_rotation, direction = _nr_roi_change_parts(mat_1_roi, mat_roi)

    if direction == "flat":
        _remove_shape(arrow_shape)
        _add_centered_text_overlay(slide, arrow_shape, "=", rgb=_GRAY)
        return

    try:
        fill = arrow_shape.fill
        fill.solid()
        fill.fore_color.rgb = change_rgb
        fill.transparency = 0.0
    except Exception:
        pass

    try:
        line = arrow_shape.line
        line.fill.solid()
        line.fill.fore_color.rgb = change_rgb
        line.fill.transparency = 0.0
    except Exception:
        pass

    if direction == "down":
        # Rotating the whole arrow shape also rotates its text, so keep the
        # arrow geometry on the shape and render the % text as a separate
        # transparent overlay in the same position.
        if getattr(arrow_shape, "has_text_frame", False):
            _set_shape_text_preserve_formatting(arrow_shape, "", rgb=_WHITE)
        try:
            arrow_shape.rotation = arrow_rotation
        except Exception:
            pass
        _add_centered_text_overlay(slide, arrow_shape, change_text, rgb=_WHITE)
        return

    # Positive case: keep the original up-arrow and write text directly into it.
    if getattr(arrow_shape, "has_text_frame", False):
        _set_shape_text_preserve_formatting(arrow_shape, change_text, rgb=_WHITE)

    try:
        arrow_shape.rotation = 0.0
    except Exception:
        pass

def _populate_media_kpis_summary_slide(
    slide,
    *,
    brand: str,
    year_range: ScopeYearRange,
    media_template_brand_df: pd.DataFrame,
) -> None:
    _apply_media_kpis_placeholders(
        slide,
        brand=brand,
        start_year=year_range.start_year,
        end_year=year_range.end_year,
    )
    roi_values = _update_net_revenue_roi_chart(
        slide,
        brand=brand,
        year_range=year_range,
        media_template_brand_df=media_template_brand_df,
    )
    _update_media_investment_chart(
        slide,
        brand=brand,
        year_range=year_range,
        media_template_brand_df=media_template_brand_df,
    )
    _update_media_incremental_volume_chart(
        slide,
        brand=brand,
        year_range=year_range,
        media_template_brand_df=media_template_brand_df,
    )
    _update_nr_roi_yoy_arrow(
        slide,
        mat_1_roi=float(roi_values.mat_1_roi or 0.0),
        mat_roi=float(roi_values.mat_roi or 0.0),
    )


def populate_media_kpis_summary_slides(
    prs: Presentation,
    *,
    media_template_brand_df: pd.DataFrame,
    scope_df: pd.DataFrame,
) -> None:
    template_slide = _find_media_kpis_template_slide(prs)
    if template_slide is None:
        return

    normalized_media_template_brand_df = _normalize_media_template_brand_df(media_template_brand_df)
    brands = _ordered_unique_brands(normalized_media_template_brand_df)
    if not brands:
        return

    year_range = _resolve_scope_year_range(scope_df)

    if len(brands) == 1:
        _populate_media_kpis_summary_slide(
            template_slide,
            brand=brands[0],
            year_range=year_range,
            media_template_brand_df=normalized_media_template_brand_df,
        )
        return

    slide_mapper = WaterfallSlideMapper()
    template_idx = list(prs.slides).index(template_slide)

    for idx, brand in enumerate(brands):
        slide = slide_mapper.duplicate_slide(prs, template_slide, insert_idx=template_idx + idx + 1)
        _populate_media_kpis_summary_slide(
            slide,
            brand=brand,
            year_range=year_range,
            media_template_brand_df=normalized_media_template_brand_df,
        )

    slide_mapper.delete_slide(prs, template_slide)
